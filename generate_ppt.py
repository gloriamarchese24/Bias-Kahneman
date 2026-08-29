import qrcode
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE_URL = "https://bias-kahneman-ms2mgk3bueryreh4vj78ae.streamlit.app"

PAGINE_IT = [
    {
        "url_path": "Macchina",
        "title": "1. Incidente Stradale (Loftus & Palmer)",
        "bias": "Falsi Ricordi & Suggestionabilità (False Memories)",
        "context": "Guarda il breve video dell'incidente stradale mostrato a lezione.\n\nSuccessivamente, inquadra il QR code per stimare la velocità dei veicoli e rispondere alle domande dal tuo smartphone."
    },
    {
        "url_path": "Malattia_Asiatica",
        "title": "2. Il Paradosso della Malattia Asiatica",
        "bias": "Effetto Framing / Formulazione della Scelta (Framing Effect)",
        "context": "Un'epidemia asiatica molto contagiosa rischia di colpire 600 persone. Sono stati proposti due diversi programmi sanitari d'intervento per affrontarla.\n\nInquadra il QR code con il tuo smartphone per scegliere il programma migliore secondo te."
    },
    {
        "url_path": "Framing_AI",
        "title": "3. Chirurgia Robotica con Intelligenza Artificiale",
        "bias": "Framing Medico (Positive vs Negative Framing)",
        "context": "Un nuovo software robotico AI deve eseguire un intervento chirurgico ad alto rischio su 100 pazienti in condizioni critiche.\n\nInquadra il QR code per leggere i dati statistici ed esprimere la tua decisione."
    },
    {
        "url_path": "Ancoraggio_Gandhi",
        "title": "4. Stima Biografica di Gandhi",
        "bias": "Effetto Ancoraggio / Valutazione di Riferimento (Anchoring Effect)",
        "context": "Mettiti alla prova con due domande sulla biografia e sulla storia di Gandhi.\n\nInquadra il QR code con il tuo smartphone per inserire la tua stima."
    },
    {
        "url_path": "Ancoraggio_Roulette",
        "title": "5. Valutazione Diagnostica Ospedaliera",
        "bias": "Ancoraggio Numerico Irrilevante (Arbitrary Anchoring)",
        "context": "Osserva il numero estratto dalla ruota della fortuna mostrata a lezione.\n\nSuccessivamente, inquadra il QR code per fornire la tua stima sulla percentuale di diagnosi errate in ospedale."
    },
    {
        "url_path": "Avversione_Perdite",
        "title": "6. Decisione Finanziaria e Rischio",
        "bias": "Avversione alle Perdite & Teoria del Prospetto (Loss Aversion)",
        "context": "Ti viene proposto uno scenario economico con un premio iniziale e due opzioni finanziarie alternative.\n\nInquadra il QR code dal tuo smartphone per scegliere l'opzione che preferisci."
    },
    {
        "url_path": "Illusione_Verita",
        "title": "7. Valutazione dell'Affermazione Scientifica",
        "bias": "Illusione di Verità & Fluidità Cognitiva (Illusion of Truth)",
        "context": "Inquadra il QR code con la fotocamera del tuo smartphone per leggere l'affermazione scientifica proposta e valutare da 1 a 10 quanto ti sembra vera."
    },
    {
        "url_path": "Euristica_Disponibilita",
        "title": "8. Inventario della Personalità (Assertività)",
        "bias": "Euristica della Disponibilità & Sforzo Cognitivo (Availability Heuristic)",
        "context": "Inquadra il QR code dal tuo smartphone per completare un breve esercizio di memoria sui tuoi comportamenti passati e valutare il tuo livello personale di assertività."
    },
    {
        "url_path": "Problema_Linda",
        "title": "9. Il Profilo di Linda (Kahneman & Tversky)",
        "bias": "Fallacia della Congiunzione & Rappresentatività (Conjunction Fallacy)",
        "context": "Linda ha 31 anni, è single, brillante, laureata in filosofia e da studentessa si è occupata di giustizia sociale e manifestazioni antinucleari.\n\nInquadra il QR code per stimare la probabilità della sua occupazione attuale."
    },
    {
        "url_path": "Effetto_Alone",
        "title": "10. Valutazione del Profilo Professionale",
        "bias": "Effetto Alone & Ordine di Presentazione (Halo Effect / Asch)",
        "context": "Inquadra il QR code per leggere la descrizione dei tratti di personalità del collega fornita dai suoi compagni di lavoro e valutare la tua impressione generale su di lui."
    },
    {
        "url_path": "Effetto_Dote_Tazza",
        "title": "11. Il Mercato della Tazza (Thaler)",
        "bias": "Effetto Dote (Endowment Effect)",
        "context": "Partecipa alla simulazione di mercato per una bellissima tazza del nostro istituto.\n\nInquadra il QR code con il tuo smartphone per scoprire il tuo ruolo e indicare il tuo prezzo."
    },
    {
        "url_path": "Effetto_Dote_AI",
        "title": "12. Mercato della Licenza Software AI",
        "bias": "Effetto Dote applicato ai Beni Digitali (Endowment Effect)",
        "context": "Partecipa alla trattativa di mercato per una rarissima licenza software diagnostica AI.\n\nInquadra il QR code per scoprire il tuo ruolo nella trattativa e la tua offerta in euro (€)."
    },
    {
        "url_path": "Costi_Sommersi_Teatro",
        "title": "13. Lo Spettacolo a Teatro (Kahneman)",
        "bias": "Fallacia dei Costi Sommersi (Sunk Cost Fallacy)",
        "context": "È la sera dello spettacolo teatrale che ti interessava molto, ma improvvisamente scoppia una tormenta di neve spaventosa.\n\nInquadra il QR code dal tuo smartphone per indicare cosa decidi di fare."
    },
    {
        "url_path": "Costi_Sommersi_AI",
        "title": "14. Gestione Progetto di Ricerca AI",
        "bias": "Costi Sommersi negli Investimenti di Sviluppo (Sunk Cost in R&D)",
        "context": "Sei a capo di un team di ricerca per un nuovo algoritmo diagnostico. Durante i lavori, Google rilascia un software gratuito tecnicamente superiore.\n\nInquadra il QR code per decidere il futuro del tuo progetto."
    },
    {
        "url_path": "Effetto_Default",
        "title": "15. Modulo di Assunzione Ospedaliera",
        "bias": "Effetto Default & Architettura delle Scelte (Default Effect / Nudge)",
        "context": "Immagina di compilare i moduli di onboarding per l'assunzione come dipendente ospedaliero.\n\nInquadra il QR code con lo smartphone per completare la tua scelta sul modulo."
    },
    {
        "url_path": "Priming_Associativo",
        "title": "16. Test di Completamento Parole",
        "bias": "Priming Semantico / Associativo (Associative Priming)",
        "context": "Inquadra il QR code per leggere velocemente le parole che appariranno sul tuo schermo e completare la parola mancante nel minor tempo possibile."
    },
    {
        "url_path": "Dunning_Kruger",
        "title": "17. Valutazione delle Abilità Accademiche",
        "bias": "Illusione di Superiorità & Effetto Dunning-Kruger (Overconfidence)",
        "context": "Inquadra il QR code con il tuo smartphone per rispondere in modo del tutto anonimo alla domanda di autovalutazione delle tue capacità accademiche."
    },
    {
        "url_path": "WYSIATI",
        "title": "18. WYSIATI — Persistenza dell'Ipotesi (Bruner & Potter)",
        "bias": "WYSIATI ('What You See Is All There Is') & Bias di Conferma Visivo",
        "context": "Guarda la sequenza visiva sul tuo smartphone dopo aver inquadrato il QR code.\n\n⚡ IMPORTANTISSIMO: Non ci sono opzioni a scelta multipla! Scrivi nello spazio sul tuo smartphone LA PRIMA COSA CHE VEDI (la tua primissima impressione non appena l'immagine compare) e invia subito!"
    },
    {
        "url_path": "Illusione_Focalizzazione",
        "title": "19. Sondaggio sul Benessere Personale",
        "bias": "Illusione di Focalizzazione (Focusing Illusion)",
        "context": "Inquadra il QR code per rispondere a un breve sondaggio sul tuo benessere generale e sulla tua soddisfazione di vita."
    },
    {
        "url_path": "Base_Rate_Neglect",
        "title": "20. Il Paradosso Diagnostico",
        "bias": "Disattenzione per la Frequenza di Base (Base Rate Neglect)",
        "context": "Una malattia rara colpisce l'1% della popolazione mondiale. Un test in grado di rilevarla è accurato al 95%. Risulti POSITIVO.\n\nInquadra il QR code per stimare la reale probabilità di avere la malattia."
    },
    {
        "url_path": "L_Esca",
        "title": "21. Scelta dell'Abbonamento (The Economist)",
        "bias": "Effetto Esca / Attrazione (Decoy Effect)",
        "context": "Devi scegliere la formula d'abbonamento alla rivista The Economist che ritieni più conveniente.\n\nInquadra il QR code con lo smartphone per selezionare la tua scelta."
    },
    {
        "url_path": "Regressione_Media",
        "title": "22. L'Effetto dell'Istruzione di Volo",
        "bias": "Illusioni Causali vs Regressione alla Media (Regression to the Mean)",
        "context": "Gli istruttori di volo israeliani notano che sgridare dopo un errore migliora la manovra successiva, mentre lodare la peggiora.\n\nInquadra il QR code per spiegare perché si verifica questo fenomeno."
    },
    {
        "url_path": "Bias_Conferma_Wason",
        "title": "23. Il Compito 2 - 4 - 6",
        "bias": "Bias di Conferma / Verifica delle Ipotesi (Wason, 1960)",
        "context": "Ti viene mostrata la sequenza numerica 2 - 4 - 6, che segue una regola segreta inventata dal docente.\n\nInquadra il QR code per scegliere la terzetta di numeri che intendi testare per scoprire la regola."
    }
]

PAGINE_EN = [
    {
        "url_path": "Macchina_EN",
        "title": "1. Car Crash Experiment (Loftus & Palmer)",
        "bias": "False Memories & Misinformation Effect",
        "context": "Watch the short video clip of the car accident shown in class.\n\nThen, scan the QR code to estimate the speed of the vehicles and answer the questions on your smartphone."
    },
    {
        "url_path": "Malattia_Asiatica_EN",
        "title": "2. The Asian Disease Problem",
        "bias": "Framing Effect (Kahneman & Tversky)",
        "context": "An unusual contagious disease is expected to kill 600 people. Two alternative healthcare programs have been proposed to combat the disease.\n\nScan the QR code with your smartphone to choose the best program according to you."
    },
    {
        "url_path": "Framing_AI_EN",
        "title": "3. Robotic AI Surgery Authorization",
        "bias": "Medical Framing Effect (Positive vs Negative)",
        "context": "A new AI robotic surgical software is set to perform a high-risk operation on 100 critical patients.\n\nScan the QR code to read the statistical data and make your decision."
    },
    {
        "url_path": "Ancoraggio_Gandhi_EN",
        "title": "4. Gandhi Age Estimation",
        "bias": "Anchoring Effect",
        "context": "Test your knowledge with two questions regarding Gandhi's life and age at death.\n\nScan the QR code with your smartphone to enter your estimates."
    },
    {
        "url_path": "Ancoraggio_Roulette_EN",
        "title": "5. Hospital Misdiagnosis Estimation",
        "bias": "Irrelevant Numerical Anchoring",
        "context": "Observe the number spun on the wheel of fortune shown in class.\n\nThen, scan the QR code to provide your estimate of the percentage of hospital misdiagnoses."
    },
    {
        "url_path": "Avversione_Perdite_EN",
        "title": "6. Financial Decision & Risk Choice",
        "bias": "Loss Aversion & Prospect Theory",
        "context": "You are presented with a financial scenario involving an initial bonus and two alternative financial options.\n\nScan the QR code on your smartphone to pick your preferred option."
    },
    {
        "url_path": "Illusione_Verita_EN",
        "title": "7. Scientific Statement Evaluation",
        "bias": "Illusion of Truth & Cognitive Ease",
        "context": "Scan the QR code with your smartphone camera to read the scientific statement presented and rate how true it seems to you (1-10)."
    },
    {
        "url_path": "Euristica_Disponibilita_EN",
        "title": "8. Personality Inventory (Assertiveness)",
        "bias": "Availability Heuristic & Subjective Ease of Recall",
        "context": "Scan the QR code from your smartphone to complete a short memory recall exercise on your past behaviors and rate your personal level of assertiveness."
    },
    {
        "url_path": "Problema_Linda_EN",
        "title": "9. Linda's Profile (Kahneman & Tversky)",
        "bias": "Conjunction Fallacy & Representativeness Heuristic",
        "context": "Linda is 31, single, outspoken, bright, philosophy major, concerned with social justice and anti-nuclear demonstrations.\n\nScan the QR code to estimate the probability of her current occupation."
    },
    {
        "url_path": "Effetto_Alone_EN",
        "title": "10. Workplace Profile Evaluation",
        "bias": "Halo Effect & Order Effect (Asch Model)",
        "context": "Scan the QR code to read the personality description of a colleague provided by his coworkers and rate your overall impression of him."
    },
    {
        "url_path": "Effetto_Dote_Tazza_EN",
        "title": "11. The University Mug Market (Thaler)",
        "bias": "Endowment Effect",
        "context": "Participate in the market simulation for an official university mug.\n\nScan the QR code with your smartphone to discover your role and submit your price."
    },
    {
        "url_path": "Effetto_Dote_AI_EN",
        "title": "12. AI Diagnostic Software License Market",
        "bias": "Endowment Effect in Digital Assets",
        "context": "Participate in the market negotiation for a rare AI diagnostic software license.\n\nScan the QR code to discover your negotiation role and enter your valuation in Euros (€)."
    },
    {
        "url_path": "Costi_Sommersi_Teatro_EN",
        "title": "13. The Theater Ticket Scenario (Kahneman)",
        "bias": "Sunk Cost Fallacy",
        "context": "It is the evening of a theater show you really want to see, but a severe snowstorm suddenly strikes.\n\nScan the QR code from your smartphone to indicate what you decide to do."
    },
    {
        "url_path": "Costi_Sommersi_AI_EN",
        "title": "14. AI Research Project Management",
        "bias": "Sunk Cost Fallacy in R&D",
        "context": "You lead an AI research team working on a new diagnostic algorithm. During development, Google releases a free superior algorithm.\n\nScan the QR code to decide the future of your project."
    },
    {
        "url_path": "Effetto_Default_EN",
        "title": "15. Hospital Employee Onboarding Form",
        "bias": "Default Effect & Choice Architecture (Nudge)",
        "context": "Imagine completing your onboarding consent forms as a new hospital employee.\n\nScan the QR code with your smartphone to make your choice on the form."
    },
    {
        "url_path": "Priming_Associativo_EN",
        "title": "16. Word Completion Test",
        "bias": "Associative / Conceptual Priming",
        "context": "Scan the QR code to quickly read the trigger words that will appear on your screen and complete the missing word as fast as possible."
    },
    {
        "url_path": "Dunning_Kruger_EN",
        "title": "17. Academic Ability Self-Assessment",
        "bias": "Illusion of Superiority & Dunning-Kruger Effect",
        "context": "Scan the QR code with your smartphone to anonymously answer a self-evaluation question about your academic abilities."
    },
    {
        "url_path": "WYSIATI_EN",
        "title": "18. WYSIATI — Hypothesis Persistence (Bruner & Potter)",
        "bias": "WYSIATI ('What You See Is All There Is') & Belief Perseverance",
        "context": "Observe the visual sequence on your smartphone after scanning the QR code.\n\n⚡ VERY IMPORTANT: No multiple choice! Type in the box on your smartphone THE VERY FIRST THING YOU SEE (your immediate first impression as soon as the image appears) and submit right away!"
    },
    {
        "url_path": "Illusione_Focalizzazione_EN",
        "title": "19. Personal Wellbeing Survey",
        "bias": "Focusing Illusion",
        "context": "Scan the QR code to complete a short survey regarding your general wellbeing and life satisfaction."
    },
    {
        "url_path": "Base_Rate_Neglect_EN",
        "title": "20. The Diagnostic Paradox",
        "bias": "Base Rate Neglect",
        "context": "A rare disease affects 1% of the population. A diagnostic test is 95% accurate. You test POSITIVE.\n\nScan the QR code to estimate the actual probability of having the condition."
    },
    {
        "url_path": "L_Esca_EN",
        "title": "21. Subscription Choice (The Economist)",
        "bias": "Decoy Effect & Asymmetric Dominance",
        "context": "You need to select a magazine subscription offer for The Economist that you find most advantageous.\n\nScan the QR code with your smartphone to submit your choice."
    },
    {
        "url_path": "Regressione_Media_EN",
        "title": "22. Flight Instruction Effect",
        "bias": "Causal Illusion vs Regression to the Mean",
        "context": "Flight instructors notice reprimanding after a bad landing leads to a better next landing, while praise leads to a worse one.\n\nScan the QR code to evaluate why this phenomenon occurs."
    },
    {
        "url_path": "Bias_Conferma_Wason_EN",
        "title": "23. The 2 - 4 - 6 Task",
        "bias": "Confirmation Bias & Hypothesis Testing (Wason, 1960)",
        "context": "You are shown the number sequence 2 - 4 - 6, which follows a secret rule created by the instructor.\n\nScan the QR code to choose which triplet of numbers you will test first to discover the rule."
    }
]

os.makedirs("qrcodes", exist_ok=True)

def build_presentation(data_list, output_filename, main_title_text, subtitle_text, banner_text, lang="IT"):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = main_title_text
    subtitle.text = subtitle_text

    # Experiment Slides
    for item in data_list:
        url_path = item["url_path"]
        url = f"{BASE_URL}/{url_path}"
        
        # Generate QR Code Image
        qr = qrcode.QRCode(version=1, box_size=10, border=4)
        qr.add_data(url)
        qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")
        img_path = f"qrcodes/{url_path}.png"
        img.save(img_path)

        # Blank slide layout
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # 1. Slide Title Box (Top Header)
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9.0), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item["title"]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(108, 99, 255)
        
        # 2. Bias Name Subheader Box
        bias_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(9.0), Inches(0.5))
        btf = bias_box.text_frame
        btf.word_wrap = True
        bp_bias = btf.paragraphs[0]
        label_bias = "🧠 Bias Cognitivo:" if lang == "IT" else "🧠 Cognitive Bias:"
        bp_bias.text = f"{label_bias} {item['bias']}"
        bp_bias.font.size = Pt(14)
        bp_bias.font.bold = True
        bp_bias.font.color.rgb = RGBColor(255, 101, 132)

        # 3. Left Side Content Box (Context / Scenario ONLY - NO Group A/B questions exposed!)
        content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.6), Inches(5.0))
        ctf = content_box.text_frame
        ctf.word_wrap = True

        p_scen = ctf.paragraphs[0]
        label_context = "📌 Istruzioni & Contesto:" if lang == "IT" else "📌 Context & Instructions:"
        p_scen.text = f"{label_context}\n{item['context']}"
        p_scen.font.size = Pt(15)
        p_scen.font.color.rgb = RGBColor(50, 50, 50)

        # 4. Right Side QR Code Image
        left_img = Inches(6.5)
        top_img = Inches(1.8)
        width_img = Inches(3.2)
        slide.shapes.add_picture(img_path, left_img, top_img, width=width_img)

        # 5. Banner / Instruction text below QR Code
        banner_box = slide.shapes.add_textbox(Inches(6.3), Inches(5.2), Inches(3.6), Inches(1.2))
        ban_tf = banner_box.text_frame
        ban_tf.word_wrap = True
        bp = ban_tf.paragraphs[0]
        bp.text = banner_text
        bp.font.size = Pt(14)
        bp.font.bold = True
        bp.font.color.rgb = RGBColor(108, 99, 255)
        bp.alignment = PP_ALIGN.CENTER

    saved = False
    count = 1
    target_name = output_filename
    while not saved:
        try:
            prs.save(target_name)
            print(f"Presentazione {target_name} generata con successo!")
            saved = True
        except PermissionError:
            count += 1
            target_name = output_filename.replace(".pptx", f"_v{count}.pptx")

# Build Italian PowerPoint
build_presentation(
    PAGINE_IT,
    "QR_Esperimenti_Bias_Domande_IT.pptx",
    "Lezione Interattiva sui Bias Cognitivi",
    "Tutti i 22 Esperimenti Kahneman\n(Inquadra i QR code per rispondere in tempo reale!)",
    "📱 Inquadra il QR code con la tua fotocamera per partecipare",
    lang="IT"
)

# Build English PowerPoint
build_presentation(
    PAGINE_EN,
    "QR_Experiments_Cognitive_Biases_EN.pptx",
    "Interactive Lesson on Cognitive Biases",
    "All 22 Kahneman Experiments\n(Scan the QR code to answer in real-time!)",
    "📱 Scan the QR code with your smartphone camera to participate",
    lang="EN"
)
